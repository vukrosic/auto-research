#!/usr/bin/env python3
"""Generate a PowerPoint overview of the autoresearch system."""

from __future__ import annotations

import argparse
import json
import re
from collections import Counter
from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as SHAPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
REPORTS = ROOT / "reports"
PROJECTS = ROOT / "projects"
EXPERIMENTS = ROOT / "experiments"
GOALS = ROOT / "goals"
STATE_MD = ROOT / "STATE.md"
ACTIVE_GOALS_MD = GOALS / "ACTIVE.md"
LOOP_RESULTS = REPORTS / "loop_results.jsonl"

TITLE_FONT = "Aptos Display"
BODY_FONT = "Aptos"
MONO_FONT = "Aptos Mono"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def rgb(hex_value: str) -> RGBColor:
    hex_value = hex_value.strip().lstrip("#")
    return RGBColor(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


BG = rgb("F6F1E8")
PAPER = rgb("FFFDF9")
SOFT = rgb("EFE2D0")
INK = rgb("1E2A32")
MUTED = rgb("57666E")
ACCENT = rgb("C55C2D")
ACCENT_DARK = rgb("8F3F1F")
GREEN = rgb("2F6F5E")
GOLD = rgb("D49A2A")
BLUE = rgb("325E73")
RED = rgb("A44747")
LINE = rgb("D4C4AF")
WHITE = rgb("FFFFFF")


STATUS_COLORS = {
    "running": BLUE,
    "pending": GOLD,
    "done": GREEN,
    "validated_winner": GREEN,
    "promoted": ACCENT,
    "rejected": RED,
    "failed": RED,
}


def utcnow() -> datetime:
    return datetime.now(timezone.utc)


def iso_now() -> str:
    return utcnow().strftime("%Y-%m-%d %H:%M UTC")


def load_json(path: Path, default: dict | list | None = None):
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        if default is None:
            return {}
        return default


def read_text(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8")
    except Exception:
        return ""


def short_name(name: str) -> str:
    if name.startswith("explore_"):
        return name[len("explore_") :]
    return name


def fmt_metric(value) -> str:
    try:
        return f"{float(value):.4f}"
    except Exception:
        return "-"


def fmt_delta(value) -> str:
    try:
        val = float(value)
    except Exception:
        return "-"
    sign = "+" if val >= 0 else ""
    return f"{sign}{val:.4f}"


def fmt_bool(flag: bool) -> str:
    return "Yes" if flag else "No"


def fit_lines(items: list[str], limit: int) -> list[str]:
    out = []
    for item in items[:limit]:
        out.append(item)
    return out


def parse_active_goals() -> list[dict]:
    lines = read_text(ACTIVE_GOALS_MD).splitlines()
    rows = []
    for line in lines:
        stripped = line.strip()
        if not stripped.startswith("|"):
            continue
        if "Goal" in stripped or "---" in stripped:
            continue
        parts = [part.strip() for part in stripped.strip("|").split("|")]
        if len(parts) < 5:
            continue
        rows.append(
            {
                "goal": parts[0],
                "mission": parts[1],
                "current_best": parts[2],
                "target": parts[3],
                "deadline": parts[4],
            }
        )
    return rows


def load_projects() -> list[dict]:
    projects = []
    for path in sorted(PROJECTS.glob("*.json")):
        payload = load_json(path, {})
        payload["_file"] = path.name
        payload["_name"] = path.stem
        projects.append(payload)
    return projects


def snapshot_records() -> list[dict]:
    records = []
    for project_dir in sorted(EXPERIMENTS.iterdir()):
        snapshots_dir = project_dir / "snapshots"
        if not snapshots_dir.is_dir():
            continue
        for snap in sorted(snapshots_dir.iterdir()):
            if not snap.is_dir():
                continue
            meta = load_json(snap / "meta.json", {})
            result = load_json(snap / "result.json", {})
            status = read_text(snap / "status").strip() or "missing"
            dispatched_at = read_text(snap / "dispatched_at").strip()
            gpu = read_text(snap / "gpu").strip()
            records.append(
                {
                    "project": project_dir.name,
                    "name": snap.name,
                    "status": status,
                    "meta": meta,
                    "result": result,
                    "gpu": gpu,
                    "dispatched_at": dispatched_at,
                    "path": snap,
                }
            )
    return records


def load_loop_results() -> list[dict]:
    rows = []
    if not LOOP_RESULTS.exists():
        return rows
    for line in LOOP_RESULTS.read_text(encoding="utf-8").splitlines():
        line = line.strip()
        if not line:
            continue
        try:
            rows.append(json.loads(line))
        except json.JSONDecodeError:
            continue
    return rows


def extract_key_finding() -> str:
    text = read_text(STATE_MD)
    match = re.search(r"\*\*key finding so far:\*\*\s*(.+)", text)
    if match:
        return match.group(1).strip()
    return "Primary records define the truth; derived state is only a view."


def build_context() -> dict:
    active_goals = parse_active_goals()
    projects = load_projects()
    snapshots = snapshot_records()
    loop_rows = load_loop_results()

    project_counts = Counter("enabled" if p.get("enabled", True) else "disabled" for p in projects)
    status_counts = Counter(record["status"] for record in snapshots)
    running = [record for record in snapshots if record["status"] == "running"]
    pending = sorted(
        [record for record in snapshots if record["status"] == "pending"],
        key=lambda record: (record["project"], record["name"]),
    )

    parameter_golf = next((p for p in projects if p.get("name") == "parameter-golf"), {})
    current_best = load_json(EXPERIMENTS / "parameter-golf" / "current_best.json", {})

    best_loop = None
    if loop_rows:
        valid = [row for row in loop_rows if row.get("val_bpb") is not None]
        if valid:
            best_loop = min(valid, key=lambda row: float(row["val_bpb"]))

    latest_loop = loop_rows[-1] if loop_rows else {}

    return {
        "generated_at": iso_now(),
        "active_goals": active_goals,
        "projects": projects,
        "project_counts": project_counts,
        "snapshots": snapshots,
        "status_counts": status_counts,
        "running": running,
        "pending": pending,
        "loop_rows": loop_rows,
        "best_loop": best_loop,
        "latest_loop": latest_loop,
        "parameter_golf": parameter_golf,
        "current_best": current_best,
        "key_finding": extract_key_finding(),
    }


def set_background(slide, color: RGBColor = BG):
    shape = slide.shapes.add_shape(SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_accent_band(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_panel(slide, left, top, width, height, fill=PAPER, line=LINE):
    shape = slide.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.2)
    return shape


def style_runs(paragraph, font_name=BODY_FONT, size=18, color=INK, bold=False):
    for run in paragraph.runs:
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    font_name=BODY_FONT,
    size=18,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    paragraph = tf.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    style_runs(paragraph, font_name=font_name, size=size, color=color, bold=bold)
    return box


def add_bullets(
    slide,
    left,
    top,
    width,
    height,
    items: list[str],
    *,
    size=16,
    color=MUTED,
    bullet="- ",
    font_name=BODY_FONT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    for idx, item in enumerate(items):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        paragraph.text = f"{bullet}{item}"
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_after = Pt(5)
        style_runs(paragraph, font_name=font_name, size=size, color=color, bold=False)
    return box


def add_card(
    slide,
    left,
    top,
    width,
    height,
    *,
    title,
    body_lines: list[str] | None = None,
    fill=PAPER,
    line=LINE,
    accent=None,
    title_size=22,
    body_size=16,
):
    add_panel(slide, left, top, width, height, fill=fill, line=line)
    if accent is not None:
        add_accent_band(slide, left, top, Inches(0.15), height, accent)
    add_textbox(
        slide,
        left + Inches(0.2),
        top + Inches(0.15),
        width - Inches(0.3),
        Inches(0.45),
        title,
        font_name=TITLE_FONT,
        size=title_size,
        color=INK,
        bold=True,
    )
    if body_lines:
        add_bullets(
            slide,
            left + Inches(0.22),
            top + Inches(0.65),
            width - Inches(0.35),
            height - Inches(0.8),
            body_lines,
            size=body_size,
            color=MUTED,
        )


def add_kpi(slide, left, top, width, height, label, value, accent):
    add_panel(slide, left, top, width, height, fill=PAPER, line=LINE)
    add_accent_band(slide, left, top, Inches(0.12), height, accent)
    add_textbox(
        slide,
        left + Inches(0.22),
        top + Inches(0.18),
        width - Inches(0.32),
        Inches(0.35),
        label.upper(),
        font_name=BODY_FONT,
        size=11,
        color=MUTED,
        bold=True,
    )
    add_textbox(
        slide,
        left + Inches(0.22),
        top + Inches(0.45),
        width - Inches(0.32),
        Inches(0.65),
        value,
        font_name=TITLE_FONT,
        size=28,
        color=INK,
        bold=True,
    )


def add_slide_title(slide, title, subtitle=None, page=None):
    add_textbox(slide, Inches(0.7), Inches(0.38), Inches(8.6), Inches(0.5), title, font_name=TITLE_FONT, size=28, color=INK, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.72), Inches(0.78), Inches(8.9), Inches(0.32), subtitle, font_name=BODY_FONT, size=13, color=MUTED)
    add_accent_band(slide, Inches(0.7), Inches(1.08), Inches(11.9), Pt(2), ACCENT)
    if page is not None:
        add_textbox(slide, Inches(12.1), Inches(0.38), Inches(0.5), Inches(0.3), str(page), font_name=BODY_FONT, size=12, color=MUTED, bold=True, align=PP_ALIGN.RIGHT)


def add_footer(slide, text):
    add_textbox(slide, Inches(0.72), Inches(7.12), Inches(5.8), Inches(0.2), text, font_name=BODY_FONT, size=9, color=MUTED)


def add_chip(slide, left, top, width, text, fill, text_color=INK):
    shape = slide.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    add_textbox(
        slide,
        left + Inches(0.08),
        top + Inches(0.05),
        width - Inches(0.16),
        Inches(0.22),
        text,
        font_name=BODY_FONT,
        size=10,
        color=text_color,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def add_flow_box(slide, left, top, width, height, title, caption, accent):
    add_panel(slide, left, top, width, height, fill=PAPER, line=LINE)
    add_accent_band(slide, left, top, Inches(0.09), height, accent)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.13), width - Inches(0.24), Inches(0.34), title, font_name=TITLE_FONT, size=18, color=INK, bold=True)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.5), width - Inches(0.24), Inches(0.6), caption, font_name=BODY_FONT, size=12, color=MUTED)


def add_chevron(slide, left, top, width, height, color=SOFT):
    shape = slide.shapes.add_shape(SHAPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def draw_status_bars(slide, counts: Counter, left, top, width):
    ordered = [
        "running",
        "pending",
        "done",
        "validated_winner",
        "rejected",
        "failed",
        "promoted",
    ]
    max_count = max([counts.get(name, 0) for name in ordered] + [1])
    label_w = Inches(1.65)
    bar_left = left + label_w
    bar_w = width - label_w - Inches(0.4)
    row_h = Inches(0.34)
    gap = Inches(0.14)
    for idx, name in enumerate(ordered):
        count = counts.get(name, 0)
        y = top + idx * (row_h + gap)
        add_textbox(slide, left, y + Inches(0.03), label_w - Inches(0.1), Inches(0.18), name.replace("_", " "), size=12, color=INK, bold=True)
        bg = slide.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, bar_left, y, bar_w, row_h)
        bg.fill.solid()
        bg.fill.fore_color.rgb = SOFT
        bg.line.fill.background()
        fill_w = max(Inches(0.14), int(bar_w * count / max_count))
        fg = slide.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, bar_left, y, fill_w, row_h)
        fg.fill.solid()
        fg.fill.fore_color.rgb = STATUS_COLORS.get(name, BLUE)
        fg.line.fill.background()
        add_textbox(slide, bar_left + bar_w - Inches(0.35), y + Inches(0.03), Inches(0.3), Inches(0.18), str(count), size=12, color=INK, bold=True, align=PP_ALIGN.RIGHT)


def draw_result_chart(slide, rows: list[dict], baseline: float, left, top, width, height):
    if not rows:
        add_card(
            slide,
            left,
            top,
            width,
            height,
            title="No loop results yet",
            body_lines=["Run `scripts/run_loop.sh` or dispatch experiments to populate the chart."],
            accent=ACCENT,
        )
        return

    rows = [row for row in rows if row.get("val_bpb") is not None]
    if not rows:
        add_card(
            slide,
            left,
            top,
            width,
            height,
            title="Loop results missing metrics",
            body_lines=["The JSONL log exists, but the current rows do not contain `val_bpb`."],
            accent=ACCENT,
        )
        return

    rows = sorted(rows, key=lambda row: float(row["val_bpb"]))[:8]
    values = [float(row["val_bpb"]) for row in rows]
    lo = min(values + [baseline]) - 0.004
    hi = max(values + [baseline]) + 0.004
    label_w = Inches(2.15)
    value_w = Inches(0.95)
    bar_left = left + label_w
    bar_w = width - label_w - value_w - Inches(0.2)
    baseline_x = bar_left + int((baseline - lo) / (hi - lo) * bar_w)

    axis = slide.shapes.add_shape(SHAPE.RECTANGLE, bar_left, top + height - Inches(0.25), bar_w, Pt(1.5))
    axis.fill.solid()
    axis.fill.fore_color.rgb = LINE
    axis.line.fill.background()

    guide = slide.shapes.add_shape(SHAPE.RECTANGLE, baseline_x, top + Inches(0.2), Pt(2), height - Inches(0.45))
    guide.fill.solid()
    guide.fill.fore_color.rgb = ACCENT_DARK
    guide.line.fill.background()
    add_textbox(slide, baseline_x - Inches(0.5), top, Inches(1.0), Inches(0.18), f"baseline {baseline:.4f}", size=10, color=ACCENT_DARK, bold=True, align=PP_ALIGN.CENTER)

    row_h = Inches(0.46)
    gap = Inches(0.12)
    for idx, row in enumerate(rows):
        name = short_name(row.get("_experiment", row.get("name", "?")))
        value = float(row["val_bpb"])
        delta = baseline - value
        y = top + Inches(0.35) + idx * (row_h + gap)
        add_textbox(slide, left, y + Inches(0.03), label_w - Inches(0.08), Inches(0.2), name, size=12, color=INK, bold=True)
        add_textbox(slide, left, y + Inches(0.23), label_w - Inches(0.08), Inches(0.18), f"delta {fmt_delta(delta)}", size=10, color=MUTED)

        bar_x = bar_left
        bar_len = max(Pt(8), int((value - lo) / (hi - lo) * bar_w))
        bar = slide.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, bar_x, y + Inches(0.03), bar_len, Inches(0.28))
        bar.fill.solid()
        bar.fill.fore_color.rgb = GREEN if value <= baseline else ACCENT
        bar.line.fill.background()

        add_textbox(slide, bar_left + bar_w + Inches(0.03), y + Inches(0.07), value_w - Inches(0.05), Inches(0.18), f"{value:.4f}", size=12, color=INK, bold=True)

    add_textbox(slide, bar_left, top + height - Inches(0.2), Inches(0.8), Inches(0.18), f"{lo:.4f}", size=10, color=MUTED)
    add_textbox(slide, bar_left + bar_w - Inches(0.8), top + height - Inches(0.2), Inches(0.8), Inches(0.18), f"{hi:.4f}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def slide_title(prs: Presentation, ctx: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_accent_band(slide, 0, 0, Inches(1.0), SLIDE_H, ACCENT)
    shape = slide.shapes.add_shape(SHAPE.OVAL, Inches(10.4), Inches(0.5), Inches(2.0), Inches(2.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SOFT
    shape.line.fill.background()
    add_textbox(slide, Inches(1.2), Inches(0.95), Inches(6.8), Inches(0.8), "Autoresearch", font_name=TITLE_FONT, size=32, color=INK, bold=True)
    add_textbox(slide, Inches(1.23), Inches(1.8), Inches(7.6), Inches(0.65), "A markdown-defined autonomous ML research lab", font_name=BODY_FONT, size=18, color=MUTED)
    add_textbox(
        slide,
        Inches(1.23),
        Inches(2.45),
        Inches(7.2),
        Inches(0.7),
        "Goals and policy live in files. Agents plan, dispatch, evaluate, promote, and update knowledge from the same filesystem state.",
        font_name=BODY_FONT,
        size=16,
        color=INK,
    )

    add_kpi(slide, Inches(1.2), Inches(4.6), Inches(2.15), Inches(1.2), "Active goals", str(len(ctx["active_goals"])), BLUE)
    add_kpi(slide, Inches(3.55), Inches(4.6), Inches(2.15), Inches(1.2), "Projects", str(len(ctx["projects"])), GREEN)
    add_kpi(slide, Inches(5.9), Inches(4.6), Inches(2.15), Inches(1.2), "Snapshots", str(len(ctx["snapshots"])), GOLD)

    add_card(
        slide,
        Inches(8.75),
        Inches(1.25),
        Inches(3.7),
        Inches(4.55),
        title="Live repo snapshot",
        body_lines=[
            f"Generated at {ctx['generated_at']}",
            f"Enabled projects: {ctx['project_counts'].get('enabled', 0)}",
            f"Running experiments: {len(ctx['running'])}",
            f"Pending experiments: {len(ctx['pending'])}",
            f"Best recent loop result: {fmt_metric(ctx['best_loop']['val_bpb']) if ctx['best_loop'] else '-'}",
        ],
        accent=ACCENT,
    )
    add_footer(slide, str(ROOT))


def slide_architecture(prs: Presentation, ctx: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_title(slide, "Two Layers, One System", "Lab governance defines authority. The execution engine turns that policy into remote GPU work.", page=2)

    add_card(
        slide,
        Inches(0.75),
        Inches(1.45),
        Inches(5.7),
        Inches(4.95),
        title="Lab operating system",
        body_lines=[
            "Human sets missions in `goals/*/MISSION.md` and policy in `lab/*.md`.",
            "AI owns plans, campaigns, week execution, knowledge updates, and reporting.",
            "Status vocabulary, promotion legality, rollback rules, and cadence come from markdown, not code.",
            "If runbook mechanics disagree with lab policy, the lab docs win.",
        ],
        accent=GREEN,
    )
    add_card(
        slide,
        Inches(6.85),
        Inches(1.45),
        Inches(5.7),
        Inches(4.95),
        title="Execution engine",
        body_lines=[
            "Projects define a universal run contract in `projects/*.json`.",
            "Experiments are full snapshots with code, metadata, status, and results.",
            "Scripts preflight, dispatch, poll, collect, parse, and optionally promote runs.",
            "Remote GPUs stay thin: most intelligence remains on the control repo.",
        ],
        accent=ACCENT,
    )
    add_chip(slide, Inches(2.2), Inches(6.65), Inches(3.3), "Boundary: policy first, mechanics second", GREEN, WHITE)
    add_chip(slide, Inches(7.55), Inches(6.65), Inches(3.7), "Primary truth lives in snapshots and project records", ACCENT, WHITE)
    add_footer(slide, "Sources: LAB.md, DESIGN.md, RESEARCH.md")


def slide_planning(prs: Presentation, ctx: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_title(slide, "Planning Hierarchy", "Plans cascade down from a human mission. Results and knowledge cascade back up.", page=3)

    boxes = [
        ("MISSION", "Human authored goal in `goals/<slug>/MISSION.md`", ACCENT),
        ("YEAR", "AI roadmap for the full mission horizon", BLUE),
        ("QUARTER", "Theme selection and major bets", BLUE),
        ("MONTH", "Execution themes and expected outcomes", BLUE),
        ("WEEK", "Concrete work list and GPU usage", GREEN),
        ("CAMPAIGN", "Research axis with waves and pivots", GREEN),
        ("QUEUE", "Goal queue or project pending set", GOLD),
        ("SNAPSHOT", "Dispatchable experiment with code + meta", ACCENT),
    ]
    left = Inches(0.7)
    top = Inches(1.55)
    box_w = Inches(1.43)
    step = Inches(1.52)
    for idx, (title, caption, accent) in enumerate(boxes):
        y = top + Inches(0.18) * idx
        add_flow_box(slide, left + step * idx, y, box_w, Inches(1.0), title, caption, accent)
        if idx < len(boxes) - 1:
            add_chevron(slide, left + step * idx + box_w + Inches(0.03), y + Inches(0.27), Inches(0.22), Inches(0.34), color=SOFT)

    goal_names = [row["goal"] for row in ctx["active_goals"]]
    add_card(
        slide,
        Inches(9.55),
        Inches(1.48),
        Inches(2.8),
        Inches(2.05),
        title="Active goals",
        body_lines=fit_lines(goal_names, 4) or ["No active goals parsed from ACTIVE.md"],
        accent=ACCENT,
        body_size=14,
    )
    add_card(
        slide,
        Inches(9.55),
        Inches(3.78),
        Inches(2.8),
        Inches(2.05),
        title="Key rule",
        body_lines=[
            "The AI can create plans, queues, campaigns, and reports.",
            "The AI cannot author or change a goal mission file.",
        ],
        accent=GREEN,
        body_size=14,
    )
    add_footer(slide, "Sources: goals/README.md, AGENTS.md, lab/13_PLANNING_HIERARCHY.md")


def slide_filesystem(prs: Presentation, ctx: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_title(slide, "The Control Plane Is the Filesystem", "Directories are not just storage. They are the operating model.", page=4)

    add_panel(slide, Inches(0.75), Inches(1.55), Inches(5.25), Inches(4.95), fill=rgb("FBF6EE"), line=LINE)
    tree = "\n".join(
        [
            "autoresearch/",
            "  lab/          permanent rules and templates",
            "  goals/        mission, plans, queues, reports",
            "  projects/     project run contracts",
            "  experiments/  base + snapshots + frontier",
            "  knowledge/    wins, failures, training notes",
            "  state/        derived dashboards",
            "  scripts/      execution engine",
            "  reports/      generated summaries",
        ]
    )
    add_textbox(slide, Inches(1.0), Inches(1.9), Inches(4.7), Inches(4.1), tree, font_name=MONO_FONT, size=16, color=INK)

    add_card(
        slide,
        Inches(6.35),
        Inches(1.55),
        Inches(5.9),
        Inches(2.0),
        title="Primary records",
        body_lines=[
            "`experiments/<project>/snapshots/*/status` is the real run state.",
            "`meta.json` and `result.json` capture the full experiment record.",
            "`current_best.json` holds the frontier and stage baselines.",
        ],
        accent=ACCENT,
    )
    add_card(
        slide,
        Inches(6.35),
        Inches(3.8),
        Inches(5.9),
        Inches(2.0),
        title="Derived views",
        body_lines=[
            "`state/*.md` files are dashboards generated from primary records.",
            "If a state file disagrees with snapshot records, regenerate the state and trust the snapshots.",
            "This makes handoffs easy without inventing a separate service layer.",
        ],
        accent=GREEN,
    )
    add_footer(slide, "Sources: state/README.md, LAB.md")


def slide_data_model(prs: Presentation, ctx: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_title(slide, "Core Data Contracts", "The system stays project-agnostic by keeping a small, explicit contract per layer.", page=5)

    add_card(
        slide,
        Inches(0.75),
        Inches(1.5),
        Inches(3.0),
        Inches(2.15),
        title="Project config",
        body_lines=[
            "Metric name and direction",
            "Run command and process pattern",
            "Log and summary parsing rules",
            "Stage budgets and thresholds",
        ],
        accent=BLUE,
    )
    add_card(
        slide,
        Inches(3.98),
        Inches(1.5),
        Inches(3.0),
        Inches(2.15),
        title="Snapshot meta",
        body_lines=[
            "Hypothesis, parent base, stage",
            "Expected duration and threshold",
            "Environment overrides",
            "Goal ownership and change summary",
        ],
        accent=GOLD,
    )
    add_card(
        slide,
        Inches(7.2),
        Inches(1.5),
        Inches(2.9),
        Inches(2.15),
        title="Result record",
        body_lines=[
            "Primary metric and optional secondary metrics",
            "Runtime, log tail, collection lag",
            "GPU, process timestamps, estimate error",
        ],
        accent=GREEN,
    )
    add_card(
        slide,
        Inches(10.32),
        Inches(1.5),
        Inches(2.0),
        Inches(2.15),
        title="State views",
        body_lines=[
            "Frontier",
            "Active runs",
            "Adjudication queue",
            "Open questions",
        ],
        accent=ACCENT,
    )

    add_card(
        slide,
        Inches(0.75),
        Inches(4.1),
        Inches(11.55),
        Inches(1.8),
        title="Status vocabulary",
        body_lines=[
            "pending -> running -> done or failed -> rejected or validated_winner -> promoted",
            "Special cases: stale_winner for old-base wins, rollback_invalidated for later reversals.",
            "These names are part of lab policy and are shared across agents and scripts.",
        ],
        accent=ACCENT,
        body_size=15,
    )
    add_footer(slide, "Sources: RESEARCH.md, lab/04_EXPERIMENT_GOVERNANCE.md, projects/*.json")


def slide_lifecycle(prs: Presentation, ctx: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_title(slide, "Lifecycle of One Idea", "A hypothesis becomes a snapshot, then compute, then a frontier decision or a learning.", page=6)

    titles = [
        ("1. Idea", "Agent identifies a promising change", BLUE),
        ("2. Snapshot", "Create full repo copy plus meta", GOLD),
        ("3. Preflight", "Check invariants, timing, goal fit", ACCENT),
        ("4. Dispatch", "Rsync to GPU and start wrapper", ACCENT),
        ("5. Run", "Remote process writes logs and runtime", BLUE),
        ("6. Collect", "Pull logs and build result.json", GREEN),
        ("7. Adjudicate", "Compare to same-step baseline", GREEN),
        ("8. Promote / learn", "Promote, rebase, or record failure", ACCENT),
    ]
    left = Inches(0.5)
    top = Inches(2.0)
    width = Inches(1.35)
    for idx, (title, caption, accent) in enumerate(titles):
        x = left + idx * Inches(1.55)
        add_flow_box(slide, x, top, width, Inches(1.45), title, caption, accent)
        if idx < len(titles) - 1:
            add_chevron(slide, x + width + Inches(0.04), top + Inches(0.48), Inches(0.22), Inches(0.42))

    add_card(
        slide,
        Inches(1.3),
        Inches(4.55),
        Inches(4.4),
        Inches(1.45),
        title="Automation hooks",
        body_lines=[
            "`new_experiment.sh`, `preflight_experiment.py`, `dispatch.sh`",
            "`check_experiment.sh`, `collect_result.sh`, `build_result.py`",
        ],
        accent=BLUE,
        body_size=14,
    )
    add_card(
        slide,
        Inches(6.15),
        Inches(4.55),
        Inches(5.1),
        Inches(1.45),
        title="Decision branch",
        body_lines=[
            "Validated latest-base wins can promote. Old-base wins become `stale_winner` and must be rebased.",
            "Everything else still feeds `knowledge/` and future planning.",
        ],
        accent=GREEN,
        body_size=14,
    )
    add_footer(slide, "Sources: DESIGN.md, RESEARCH.md")


def slide_execution_modes(prs: Presentation, ctx: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_title(slide, "Execution Modes and Scheduling", "The repo supports interactive cycles, batch loops, and deadline-bound goal windows.", page=7)

    add_card(
        slide,
        Inches(0.75),
        Inches(1.6),
        Inches(3.75),
        Inches(3.75),
        title="Interactive cycle",
        body_lines=[
            "`python3 scripts/autonomous_lab.py cycle`",
            "Multi-project aware and file-lock protected.",
            "Reconciles state, checks running work, adjudicates completed runs, and dispatches the next eligible snapshot.",
        ],
        accent=BLUE,
    )
    add_card(
        slide,
        Inches(4.78),
        Inches(1.6),
        Inches(3.75),
        Inches(3.75),
        title="Batch loop",
        body_lines=[
            "`bash scripts/run_loop.sh <N>`",
            "Single-project serial queue for deterministic, low-touch screening.",
            "Appends a machine-readable JSONL record after each completed experiment.",
        ],
        accent=ACCENT,
    )
    add_card(
        slide,
        Inches(8.8),
        Inches(1.6),
        Inches(3.75),
        Inches(3.75),
        title="Goal window",
        body_lines=[
            "`bash scripts/run_goal_window.sh <goal>`",
            "Materializes goal-specific queue entries, runs cycles to deadline, then emits a deadline report.",
            "Short sprints can use training windows tied to first dispatch rather than planning start time.",
        ],
        accent=GREEN,
    )

    chips = [
        ("Goal queue overrides legacy queue", BLUE),
        ("Preflight is a hard gate", ACCENT),
        ("One GPU, one goal at a time", GOLD),
        ("Measured runtime beats guesses", GREEN),
    ]
    x = Inches(0.95)
    for label, color in chips:
        width = Inches(2.6) if len(label) < 26 else Inches(3.1)
        add_chip(slide, x, Inches(5.85), width, label, color, WHITE)
        x += width + Inches(0.2)
    add_footer(slide, "Sources: scripts/run_loop.sh, scripts/run_goal_window.sh, scripts/autonomous_lab.py")


def slide_live_snapshot(prs: Presentation, ctx: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_title(slide, "Live Repository Snapshot", f"Primary records read at render time: {ctx['generated_at']}", page=8)

    add_kpi(slide, Inches(0.75), Inches(1.45), Inches(2.3), Inches(1.1), "Active goals", str(len(ctx["active_goals"])), BLUE)
    add_kpi(slide, Inches(3.25), Inches(1.45), Inches(2.3), Inches(1.1), "Enabled projects", str(ctx["project_counts"].get("enabled", 0)), GREEN)
    add_kpi(slide, Inches(5.75), Inches(1.45), Inches(2.3), Inches(1.1), "Total snapshots", str(len(ctx["snapshots"])), GOLD)
    add_kpi(slide, Inches(8.25), Inches(1.45), Inches(2.3), Inches(1.1), "Running now", str(len(ctx["running"])), ACCENT)

    add_card(
        slide,
        Inches(0.75),
        Inches(2.9),
        Inches(5.0),
        Inches(3.0),
        title="Current run",
        body_lines=(
            [
                f"{record['project']} / {record['name']}",
                f"stage={record['meta'].get('stage', '-')}, steps={record['meta'].get('steps', '-')}",
                f"gpu={record['gpu'] or '-'}, dispatched={record['dispatched_at'] or '-'}",
            ]
            for record in ctx["running"][:1]
        ).__next__() if ctx["running"] else ["No experiments are marked `running` right now."],
        accent=ACCENT,
    )

    add_panel(slide, Inches(6.0), Inches(2.9), Inches(3.05), Inches(3.0), fill=PAPER, line=LINE)
    add_textbox(slide, Inches(6.22), Inches(3.08), Inches(2.5), Inches(0.35), "Snapshot status mix", font_name=TITLE_FONT, size=20, color=INK, bold=True)
    draw_status_bars(slide, ctx["status_counts"], Inches(6.18), Inches(3.52), Inches(2.6))

    pending_lines = [
        f"{idx}. {record['name']}"
        for idx, record in enumerate(ctx["pending"][:6], 1)
    ]
    add_card(
        slide,
        Inches(9.25),
        Inches(2.9),
        Inches(3.05),
        Inches(3.0),
        title="Queue head",
        body_lines=pending_lines or ["No pending experiments."],
        accent=BLUE,
        body_size=14,
    )
    add_footer(slide, "Live counts come from experiments/*/snapshots/*, not from state/*.md")


def slide_parameter_golf(prs: Presentation, ctx: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_title(slide, "Parameter Golf Example Project", "The current production lane for the lab.", page=9)

    project = ctx["parameter_golf"]
    best = ctx["current_best"]
    baseline = best.get("stage_baselines", {}).get("explore", {})
    target = project.get("target")
    current_best_quant = project.get("current_best")
    gap = float(current_best_quant) - float(target) if target is not None and current_best_quant is not None else None

    add_card(
        slide,
        Inches(0.75),
        Inches(1.55),
        Inches(4.3),
        Inches(4.75),
        title="Challenge profile",
        body_lines=[
            project.get("description", "No description available."),
            f"Primary metric: {project.get('metric', '-')}, direction={project.get('metric_direction', '-')}",
            f"Current frontier: {project.get('current_best_description', '-')}",
            f"Target: {fmt_metric(target)} ; current best: {fmt_metric(current_best_quant)} ; gap: {fmt_metric(gap) if gap is not None else '-'}",
        ],
        accent=ACCENT,
    )

    stages = project.get("stages", {})
    stage_cards = [
        ("Explore", stages.get("explore", {}), BLUE),
        ("Validate", stages.get("validate", {}), GREEN),
        ("Full", stages.get("full", {}), GOLD),
    ]
    x = Inches(5.35)
    for title, payload, accent in stage_cards:
        add_card(
            slide,
            x,
            Inches(1.55),
            Inches(2.2),
            Inches(2.2),
            title=title,
            body_lines=[
                f"steps={payload.get('steps', '-')}",
                f"threshold={payload.get('threshold', '-')}",
                payload.get("description", "-"),
            ],
            accent=accent,
            body_size=13,
        )
        x += Inches(2.35)

    add_card(
        slide,
        Inches(5.35),
        Inches(4.05),
        Inches(6.95),
        Inches(2.25),
        title="Current calibrated baseline",
        body_lines=[
            f"Experiment: {best.get('experiment_name', '-')}",
            f"Explore baseline at 500 steps: val_bpb={fmt_metric(baseline.get('val_bpb'))}",
            f"Quant baseline: val_bpb_quant={fmt_metric(baseline.get('val_bpb_quant'))}",
            "Every experiment is judged against a baseline at the same step count.",
        ],
        accent=GREEN,
    )
    add_footer(slide, "Sources: projects/parameter-golf.json, experiments/parameter-golf/current_best.json")


def slide_results(prs: Presentation, ctx: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_title(slide, "Recent Explore Results", "Loop results from `reports/loop_results.jsonl`, plotted against the current 500-step baseline.", page=10)

    baseline = float(ctx["current_best"].get("stage_baselines", {}).get("explore", {}).get("val_bpb", 1.6673))
    draw_result_chart(slide, ctx["loop_rows"], baseline, Inches(0.72), Inches(1.55), Inches(8.2), Inches(5.2))

    best_row = ctx["best_loop"]
    latest_row = ctx["latest_loop"]
    summary_lines = [
        f"Best recent result: {short_name(best_row.get('_experiment', '-'))} at {fmt_metric(best_row.get('val_bpb'))}" if best_row else "Best recent result: -",
        f"Latest collected result: {short_name(latest_row.get('_experiment', '-'))} at {fmt_metric(latest_row.get('val_bpb'))}" if latest_row else "Latest collected result: -",
        ctx["key_finding"],
    ]
    add_card(
        slide,
        Inches(9.2),
        Inches(1.55),
        Inches(3.15),
        Inches(3.15),
        title="Takeaway",
        body_lines=summary_lines,
        accent=ACCENT,
        body_size=14,
    )

    top_rows = sorted([row for row in ctx["loop_rows"] if row.get("val_bpb") is not None], key=lambda row: float(row["val_bpb"]))[:3]
    leaderboard = [
        f"{idx}. {short_name(row.get('_experiment', '-'))} -> {fmt_metric(row.get('val_bpb'))}"
        for idx, row in enumerate(top_rows, 1)
    ]
    add_card(
        slide,
        Inches(9.2),
        Inches(4.95),
        Inches(3.15),
        Inches(1.8),
        title="Top recent runs",
        body_lines=leaderboard or ["No result rows found."],
        accent=GREEN,
        body_size=14,
    )
    add_footer(slide, "Metric is lower-is-better val_bpb. Green bars beat baseline; orange bars miss it.")


def slide_design_choices(prs: Presentation, ctx: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_title(slide, "Intentional Design Choices", "The system trades service complexity for explicit files, simple scripts, and auditable runs.", page=11)

    add_card(
        slide,
        Inches(0.75),
        Inches(1.55),
        Inches(3.75),
        Inches(2.15),
        title="Snapshots, not branches",
        body_lines=[
            "Full directory copies avoid branch merge churn across many experiments.",
            "Each run is self-contained and reproducible on its own.",
        ],
        accent=ACCENT,
    )
    add_card(
        slide,
        Inches(4.78),
        Inches(1.55),
        Inches(3.75),
        Inches(2.15),
        title="Markdown, not a service DB",
        body_lines=[
            "Policy, plans, and state are inspectable with a plain text editor.",
            "A new agent can recover context by reading files, not by querying hidden state.",
        ],
        accent=GREEN,
    )
    add_card(
        slide,
        Inches(8.8),
        Inches(1.55),
        Inches(3.45),
        Inches(2.15),
        title="Thin GPU side",
        body_lines=[
            "Dispatch uses SSH, rsync, nohup, and a tiny runtime wrapper.",
            "The repo remains the brain; GPUs stay replaceable workers.",
        ],
        accent=BLUE,
    )
    add_card(
        slide,
        Inches(0.75),
        Inches(4.05),
        Inches(5.55),
        Inches(2.0),
        title="What this buys you",
        body_lines=[
            "Low operational overhead, explicit handoffs, project-agnostic contracts, and fast debugging when a run goes wrong.",
            "Most failure modes remain visible in the same snapshot that requested the compute.",
        ],
        accent=GREEN,
        body_size=15,
    )
    add_card(
        slide,
        Inches(6.6),
        Inches(4.05),
        Inches(5.65),
        Inches(2.0),
        title="Known tradeoffs",
        body_lines=[
            "Snapshot copies cost disk. State dashboards can drift if not regenerated. The current deployment is mostly one-GPU serial.",
            "Literature review automation and broader multi-project onboarding are still open items.",
        ],
        accent=ACCENT,
        body_size=15,
    )
    add_footer(slide, "Sources: DESIGN.md, STATE.md")


def slide_roadmap(prs: Presentation, ctx: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_title(slide, "Next Moves", "The current repo already works. These are the highest-leverage ways to deepen it.", page=12)

    add_card(
        slide,
        Inches(0.75),
        Inches(1.55),
        Inches(3.8),
        Inches(4.7),
        title="Research roadmap",
        body_lines=[
            "Move beyond knob sweeps into code-level bets such as quantization, vocab changes, XSA, EMA/SWA, and TTT.",
            "Keep short explore lanes for signal, then validate only the best survivors.",
        ],
        accent=ACCENT,
    )
    add_card(
        slide,
        Inches(4.8),
        Inches(1.55),
        Inches(3.8),
        Inches(4.7),
        title="Systems roadmap",
        body_lines=[
            "Regenerate state views automatically at session open or after each cycle.",
            "Improve multi-GPU prioritization and project onboarding templates.",
            "Deduplicate snapshot storage when repositories get large.",
        ],
        accent=GREEN,
    )
    add_card(
        slide,
        Inches(8.85),
        Inches(1.55),
        Inches(3.45),
        Inches(4.7),
        title="Operator quickstart",
        body_lines=[
            "Read `AGENTS.md`, then `state/NOW.md`.",
            "Run `python3 scripts/autonomous_lab.py cycle` for a full agent cycle.",
            "Use `bash scripts/run_loop.sh <N>` for serial explore screening.",
            "Use `bash scripts/run_goal_window.sh <goal>` for a deadline sprint.",
        ],
        accent=BLUE,
    )
    add_footer(slide, "Generated from live repo state by scripts/make_autoresearch_system_presentation.py")


def build_presentation(ctx: dict) -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    props = prs.core_properties
    props.title = "Autoresearch System Overview"
    props.subject = "Autonomous research system architecture"
    props.author = "OpenAI Codex"
    props.keywords = "autoresearch,powerpoint,architecture,ml research"

    slide_title(prs, ctx)
    slide_architecture(prs, ctx)
    slide_planning(prs, ctx)
    slide_filesystem(prs, ctx)
    slide_data_model(prs, ctx)
    slide_lifecycle(prs, ctx)
    slide_execution_modes(prs, ctx)
    slide_live_snapshot(prs, ctx)
    slide_parameter_golf(prs, ctx)
    slide_results(prs, ctx)
    slide_design_choices(prs, ctx)
    slide_roadmap(prs, ctx)
    return prs


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--output",
        type=Path,
        default=REPORTS / f"autoresearch_system_overview_{utcnow().strftime('%Y-%m-%d')}.pptx",
        help="Output .pptx path",
    )
    args = parser.parse_args()

    REPORTS.mkdir(parents=True, exist_ok=True)
    ctx = build_context()
    prs = build_presentation(ctx)
    args.output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.output)
    print(args.output)


if __name__ == "__main__":
    main()
